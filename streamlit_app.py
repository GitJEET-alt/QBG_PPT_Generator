import base64
import copy
import io
import os
import re
import statistics
import zipfile
import tempfile

import streamlit as st
import pandas as pd
import numpy as np
from PIL import Image

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ---------- CONFIG ----------
ANCHOR_SHAPE_NAME = "SCREENSHOT_BOX"

WHITE_THRESHOLD = 240
CROP_PAD_PX = 20

EMU_PER_INCH = 914400

INK_ROW_THRESHOLD = 180
MIN_BAND_PX = 3
MAX_BAND_PX = 60

# Target height (inches) for a typical text line, measured directly from
# each image's own pixel content (see estimate_line_height_px). Used to
# derive one fixed pixel-to-EMU scale applied to every image in a batch,
# instead of stretching each one to fill a fixed box or inferring font
# size indirectly from overall image width/height.
TARGET_LINE_HEIGHT_INCHES = 0.25
# ---------------------------


st.set_page_config(page_title="PPT Generator", layout="centered")


# ---------- AUTH GATE ----------
def require_login_and_allowlist():
    u = getattr(st, "user", None)

    if not u or not u.is_logged_in:
        st.header("Private tool")
        st.write("Please sign in with your Google Workspace account.")
        if st.button("Log in with Google"):
            st.login()
        st.stop()

    email = (u.get("email") or "").strip().lower()

    raw = os.environ.get("ALLOWED_EMAILS", "")
    allowed = {e.strip().lower() for e in raw.replace(",", "\n").splitlines() if e.strip()}

    if allowed and email not in allowed:
        st.error(f"Access denied for: {email}")
        if st.button("Log out"):
            st.logout()
        st.stop()

    st.caption(f"Signed in as: {email}")


require_login_and_allowlist()
# -----------------------------


# --- UI CSS ---
st.markdown(
    """
    <style>
      .block-container {
        padding-top: 2.2rem !important;
        padding-bottom: 4.8rem !important;
        max-width: 820px !important;
      }

      div[data-testid="stVerticalBlock"] > div {
        gap: 0.8rem !important;
      }

      [data-testid="stAppViewContainer"] {
        scrollbar-width: none;
      }

      [data-testid="stAppViewContainer"]::-webkit-scrollbar {
        display: none;
      }

      footer {visibility: hidden;}
      #MainMenu {visibility: hidden;}
    </style>
    """,
    unsafe_allow_html=True,
)

st.title("📊 PPT Generator Tool")
st.write("Upload a ZIP (images) and an XLSX (mapping) to generate a PPTX.")

upload_mode = st.radio(
    "File upload mode",
    ["Separate ZIP + XLSX", "Combined ZIP (images + XLSX bundled)"],
    horizontal=True,
)

zip_file = None
xlsx_file = None
combined_file = None

if upload_mode == "Separate ZIP + XLSX":
    zip_file = st.file_uploader("Upload ZIP file", type=["zip"])
    xlsx_file = st.file_uploader("Upload XLSX file", type=["xlsx"])
else:
    combined_file = st.file_uploader(
        "Upload combined ZIP (must contain one XLSX and one images ZIP)",
        type=["zip"],
    )

include_solutions = st.checkbox("Include Solution Images")

custom_color_enabled = st.checkbox("Custom text color")
text_color_hex = None

if custom_color_enabled:
    text_color_hex = st.color_picker("Text color", value="#000000")
    st.caption("Tints dark/grayscale pixels (text) toward this color. Colorful diagrams are left alone.")

custom_template_file = st.file_uploader(
    "Custom PPT template (optional)",
    type=["pptx"],
    help=f"Must have at least one slide. Question images are placed using the default "
         f"template's '{ANCHOR_SHAPE_NAME}' position, scaled to this template's slide size "
         f"-- no anchor shape needed in the custom file itself.",
)

output_name_input = st.text_input(
    "Output file name (optional)",
    placeholder="Leave blank to use the ZIP file's name",
)


def remove_white(img: Image.Image, thr=WHITE_THRESHOLD) -> Image.Image:
    """Make near-white pixels transparent."""
    rgba = np.array(img.convert("RGBA"))
    rgb = rgba[:, :, :3]
    mask = (rgb >= thr).all(axis=2)
    rgba[mask, 3] = 0
    return Image.fromarray(rgba, "RGBA")


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))


def recolor_text(cleaned_rgba, hex_color, saturation_threshold=30):
    """
    Tint dark/grayscale ("ink") pixels toward hex_color, blending by how
    dark each pixel originally was so anti-aliased edges stay smooth.
    Colorful pixels (diagrams, highlights) are left untouched so this
    only affects text and black-and-white line art.
    """
    rgba = np.array(cleaned_rgba).astype(float)
    rgb = rgba[:, :, :3]

    max_c = rgb.max(axis=2)
    min_c = rgb.min(axis=2)
    is_ink = (max_c - min_c) <= saturation_threshold

    darkness = 1 - (max_c / 255.0)
    target = np.array(hex_to_rgb(hex_color), dtype=float)

    blended = 255 * (1 - darkness[..., None]) + target[None, None, :] * darkness[..., None]
    rgb = np.where(is_ink[..., None], blended, rgb)

    rgba[:, :, :3] = np.clip(rgb, 0, 255)
    return Image.fromarray(rgba.astype(np.uint8), "RGBA")


def crop_to_content_rgba(img_rgba: Image.Image, pad_px=CROP_PAD_PX) -> Image.Image:
    """Crop image to non-transparent content using alpha channel bounding box."""
    alpha = img_rgba.split()[-1]
    bbox = alpha.getbbox()

    if not bbox:
        return img_rgba

    x0, y0, x1, y1 = bbox
    x0 = max(0, x0 - pad_px)
    y0 = max(0, y0 - pad_px)
    x1 = min(img_rgba.width, x1 + pad_px)
    y1 = min(img_rgba.height, y1 + pad_px)

    return img_rgba.crop((x0, y0, x1, y1))


def extract_filename(cell):
    """Extract image filename from URLs/paths/cells."""
    if pd.isna(cell):
        return ""

    s = str(cell).strip()
    m = re.search(r"([^/\s]+?\.(png|jpg|jpeg|webp|bmp))", s, re.I)
    return m.group(1) if m else s


def sanitize_filename(name):
    """Strip characters that aren't valid in a Windows/Mac/Linux file name."""
    return re.sub(r'[\\/:*?"<>|]+', "_", name).strip()


def iter_all_shapes(shapes):
    for shp in shapes:
        yield shp
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_all_shapes(shp.shapes)


def find_anchor_on_slide(slide, name):
    for shp in iter_all_shapes(slide.shapes):
        if getattr(shp, "name", "").strip() == name:
            return shp
    return None


def compute_fixed_scale_size(width_px, height_px, scale, max_w, max_h):
    """
    Size an image at a fixed EMU-per-pixel scale (same scale for every
    image in the batch), shrinking only if it would overflow the
    available space -- never magnifying small images to fill it.
    """
    w = width_px * scale
    h = height_px * scale

    if w > max_w or h > max_h:
        shrink = min(max_w / w, max_h / h)
        w *= shrink
        h *= shrink

    return int(w), int(h)


def safe_extract(zipf: zipfile.ZipFile, dest_dir: str):
    """Safely extract ZIP to prevent path traversal."""
    dest_abs = os.path.abspath(dest_dir)

    for member in zipf.infolist():
        if member.is_dir():
            continue

        member_path = member.filename.replace("\\", "/").lstrip("/")
        out_path = os.path.abspath(os.path.join(dest_dir, member_path))

        if not out_path.startswith(dest_abs + os.sep):
            raise ValueError(f"Unsafe zip entry blocked: {member.filename}")

        os.makedirs(os.path.dirname(out_path), exist_ok=True)

        with zipf.open(member) as src, open(out_path, "wb") as dst:
            dst.write(src.read())


def split_combined_zip(combined_bytes):
    """Pull the XLSX mapping and the images ZIP out of a combined bundle."""
    with zipfile.ZipFile(io.BytesIO(combined_bytes), "r") as zf:
        xlsx_name = next((n for n in zf.namelist() if n.lower().endswith(".xlsx")), None)
        images_name = next((n for n in zf.namelist() if n.lower().endswith(".zip")), None)

        if not xlsx_name:
            raise RuntimeError("Combined ZIP must contain an XLSX mapping file.")

        if not images_name:
            raise RuntimeError("Combined ZIP must contain an images ZIP file.")

        return zf.read(xlsx_name), zf.read(images_name)


def build_basename_index(root_dir: str):
    """Index files by basename lowercase."""
    idx = {}

    for root, _, files in os.walk(root_dir):
        for f in files:
            idx[f.lower()] = os.path.join(root, f)

    return idx


def find_src_path(idx, src_basename):
    """Find by exact basename; fallback fuzzy by stem ignoring spaces."""
    if not src_basename:
        return None

    src_basename = str(src_basename).strip()

    if src_basename.lower() in ["nan", "none", ""]:
        return None

    p = idx.get(src_basename.lower())

    if p:
        return p

    stem = os.path.splitext(src_basename)[0].lower().replace(" ", "")

    for name, pth in idx.items():
        if os.path.splitext(name)[0].replace(" ", "") == stem:
            return pth

    return None


def estimate_line_height_px(cleaned_rgba):
    """
    Estimate a typical text-line pixel height by scanning for horizontal
    bands of dark ("ink") pixels, separated by near-blank rows. This reads
    the actual rendered font size from the image content, rather than
    inferring it indirectly from the image's overall width or height.
    """
    gray = np.array(cleaned_rgba.convert("L"))
    ink_rows = (gray < INK_ROW_THRESHOLD).sum(axis=1) > 2

    bands = []
    in_band = False
    start = 0

    for y, has_ink in enumerate(ink_rows):
        if has_ink and not in_band:
            in_band = True
            start = y
        elif not has_ink and in_band:
            in_band = False
            bands.append(y - start)

    if in_band:
        bands.append(len(ink_rows) - start)

    bands = [b for b in bands if MIN_BAND_PX <= b <= MAX_BAND_PX]

    return statistics.median(bands) if bands else None


def clean_image(image_path):
    """Load, remove white background, and crop to content."""
    with Image.open(image_path) as im:
        cleaned = remove_white(im, thr=WHITE_THRESHOLD)
    return crop_to_content_rgba(cleaned, pad_px=CROP_PAD_PX)


def place_cleaned_image(slide, cleaned, anchor_left, anchor_top, slide_w, slide_h, scale, text_color=None):
    """Place an already-cleaned image at a fixed scale, capped to fit the slide."""
    w_px, h_px = cleaned.size

    if h_px == 0:
        return False

    max_w = max(1, slide_w - anchor_left)
    max_h = max(1, slide_h - anchor_top)

    w_emu, h_emu = compute_fixed_scale_size(w_px, h_px, scale, max_w, max_h)

    to_save = recolor_text(cleaned, text_color) if text_color else cleaned

    buf = io.BytesIO()
    to_save.save(buf, format="PNG")
    buf.seek(0)

    slide.shapes.add_picture(
        buf,
        anchor_left,
        anchor_top,
        width=w_emu,
        height=h_emu,
    )

    return True


if st.button("🚀 Generate PPT"):
    if upload_mode == "Separate ZIP + XLSX":
        if not zip_file or not xlsx_file:
            st.error("Please upload both ZIP and XLSX files.")
            st.stop()
    else:
        if not combined_file:
            st.error("Please upload the combined ZIP file.")
            st.stop()

    try:
        with st.spinner("Processing..."):
            if upload_mode == "Separate ZIP + XLSX":
                xlsx_bytes = xlsx_file.getvalue()
                images_zip_bytes = zip_file.getvalue()
                default_name_source = zip_file.name
            else:
                xlsx_bytes, images_zip_bytes = split_combined_zip(combined_file.getvalue())
                default_name_source = combined_file.name

            # Read mapping
            df_raw = pd.read_excel(io.BytesIO(xlsx_bytes))

            # Exact fixed column headers
            col_order = "Display Order*"
            col_image = "Question Image"
            col_solution = "Sol Image"

            required_cols = [col_order, col_image]

            if include_solutions:
                required_cols.append(col_solution)

            missing_cols = [c for c in required_cols if c not in df_raw.columns]

            if missing_cols:
                raise RuntimeError(
                    f"Missing required column(s): {missing_cols}. "
                    f"Found columns: {list(df_raw.columns)}"
                )

            needed_cols = [col_order, col_image]

            if include_solutions:
                needed_cols.append(col_solution)

            df = df_raw[needed_cols].copy()

            rename_map = {
                col_order: "display_order",
                col_image: "question_image",
            }

            if include_solutions:
                rename_map[col_solution] = "solution_image"

            df = df.rename(columns=rename_map)

            # Keep original behavior for questions
            df = df.dropna(subset=["display_order", "question_image"]).copy()

            df["display_order"] = df["display_order"].apply(
                lambda x: int(str(x).strip().split(".")[0])
            )

            df["question_src_name"] = df["question_image"].apply(extract_filename)

            if include_solutions:
                # Do not drop question rows if solution image is missing
                df["solution_src_name"] = df["solution_image"].apply(extract_filename)

            df = df.sort_values("display_order").reset_index(drop=True)

            # Load the default template to read the anchor's position as a
            # percentage of the slide. A custom template reuses this same
            # relative position, so it doesn't need its own anchor shape.
            default_template_path = "template.pptx"

            if not os.path.exists(default_template_path):
                raise FileNotFoundError("template.pptx not found in app directory.")

            default_prs = Presentation(default_template_path)
            default_anchor = find_anchor_on_slide(default_prs.slides[0], ANCHOR_SHAPE_NAME)

            if default_anchor is None:
                raise RuntimeError(
                    f"Anchor '{ANCHOR_SHAPE_NAME}' not found on the default template's "
                    f"slide 1. Rename via Selection Pane."
                )

            anchor_left_pct = default_anchor.left / default_prs.slide_width
            anchor_top_pct = default_anchor.top / default_prs.slide_height

            if custom_template_file:
                prs = Presentation(io.BytesIO(custom_template_file.getvalue()))

                if len(prs.slides) == 0:
                    raise RuntimeError("Custom template must have at least one slide.")
            else:
                prs = default_prs

            # Use layout of slide 1
            template_layout = prs.slides[0].slide_layout

            # Snapshot slide 1's actual shapes (background, borders, logo,
            # etc.) so every newly created slide matches it exactly.
            # add_slide() on its own instead populates a new slide with the
            # layout's default placeholder shapes (e.g. "Click to add
            # title"), even ones that were deleted from slide 1 itself.
            source_slide_shapes = [copy.deepcopy(shp._element) for shp in prs.slides[0].shapes]

            def ensure_slide(i0):
                while len(prs.slides) <= i0:
                    new_slide = prs.slides.add_slide(template_layout)

                    for shp in list(new_slide.shapes):
                        new_slide.shapes._spTree.remove(shp._element)

                    for el in source_slide_shapes:
                        new_slide.shapes._spTree.append(copy.deepcopy(el))

                return prs.slides[i0]

            slide_w = prs.slide_width
            slide_h = prs.slide_height
            ANCHOR_LEFT = int(slide_w * anchor_left_pct)
            ANCHOR_TOP = int(slide_h * anchor_top_pct)

            with tempfile.TemporaryDirectory(prefix="pptgen_") as tmpdir:
                src_dir = os.path.join(tmpdir, "src")
                os.makedirs(src_dir, exist_ok=True)

                # Unzip safely
                with zipfile.ZipFile(io.BytesIO(images_zip_bytes), "r") as zf:
                    safe_extract(zf, src_dir)

                idx = build_basename_index(src_dir)

                question_missing = 0
                solution_missing = 0
                question_placed = 0
                solution_placed = 0

                # Resolve source paths and clean every distinct image once,
                # so the placement scale can be calibrated from the whole batch
                # instead of forcing each image to fill the same box area.
                resolved_rows = []
                cleaned_cache = {}

                for row_position, row in df.iterrows():
                    question_src_path = find_src_path(idx, row["question_src_name"])
                    solution_src_path = None

                    if include_solutions:
                        solution_src_path = find_src_path(idx, row.get("solution_src_name", ""))

                    resolved_rows.append((row_position, row, question_src_path, solution_src_path))

                    for src_path in (question_src_path, solution_src_path):
                        if src_path and src_path not in cleaned_cache:
                            cleaned_cache[src_path] = clean_image(src_path)

                line_heights = [estimate_line_height_px(img) for img in cleaned_cache.values()]
                line_heights = [h for h in line_heights if h]
                reference_line_height = statistics.median(line_heights) if line_heights else 20
                scale = (TARGET_LINE_HEIGHT_INCHES * EMU_PER_INCH) / reference_line_height

                for row_position, row, question_src_path, solution_src_path in resolved_rows:
                    if include_solutions:
                        question_slide_index = row_position * 2
                    else:
                        question_slide_index = int(row["display_order"]) - 1

                    question_slide = ensure_slide(question_slide_index)

                    if question_src_path:
                        ok = place_cleaned_image(
                            question_slide,
                            cleaned_cache[question_src_path],
                            ANCHOR_LEFT,
                            ANCHOR_TOP,
                            slide_w,
                            slide_h,
                            scale,
                            text_color_hex,
                        )

                        if ok:
                            question_placed += 1
                    else:
                        question_missing += 1

                    if include_solutions:
                        solution_slide_index = question_slide_index + 1
                        solution_slide = ensure_slide(solution_slide_index)

                        if solution_src_path:
                            ok = place_cleaned_image(
                                solution_slide,
                                cleaned_cache[solution_src_path],
                                ANCHOR_LEFT,
                                ANCHOR_TOP,
                                slide_w,
                                slide_h,
                                scale,
                                text_color_hex,
                            )

                            if ok:
                                solution_placed += 1
                        else:
                            solution_missing += 1

            # Save PPTX to memory
            out = io.BytesIO()
            prs.save(out)
            out.seek(0)

        if include_solutions:
            st.success(
                f"✅ Questions placed: {question_placed}. "
                f"Solutions placed: {solution_placed}. "
                f"Question missing: {question_missing}. "
                f"Solution missing: {solution_missing}."
            )
        else:
            st.success(
                f"✅ Placed {question_placed} question image(s). "
                f"Missing: {question_missing}."
            )

        custom_name = sanitize_filename(output_name_input)
        base_name = custom_name or os.path.splitext(default_name_source)[0]

        if not base_name.lower().endswith(".pptx"):
            base_name += ".pptx"

        b64_pptx = base64.b64encode(out.getvalue()).decode()

        st.components.v1.html(
            f"""
            <a id="auto_download_link"
               href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_pptx}"
               download="{base_name}"></a>
            <script>
                document.getElementById("auto_download_link").click();
            </script>
            """,
            height=0,
        )

    except Exception as e:
        st.error(f"Error: {e}")


# ---- Footer ----
st.markdown(
    """
    <div class="neet-footer">
      Powered by NEET ARD Team | academics.innovation@pw.live
    </div>

    <style>
      .neet-footer{
        position: fixed;
        left: 0;
        bottom: 0;
        width: 100%;
        background: #f3f4f6;
        border-top: 1px solid #e5e7eb;
        padding: 10px 12px;
        text-align: center;
        font-size: 14px;
        color: #4b5563;
        z-index: 9999;
      }
    </style>
    """,
    unsafe_allow_html=True,
)